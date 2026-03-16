#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
"""PowerPoint (.pptx) file parser for the RAG pipeline.

Extracts text content from PowerPoint presentations, handling text frames,
tables, group shapes, and bulleted lists. Shapes are sorted by position
(top-to-bottom, left-to-right) to preserve reading order.
"""

import logging
from io import BytesIO
from pptx import Presentation


class RAGFlowPptParser:
    """Parser for PowerPoint (.pptx) files.

    Extracts text from slides while preserving reading order by sorting shapes
    spatially. Handles text frames, tables (shape_type 19), group shapes
    (shape_type 6), and bulleted/numbered lists.

    Attributes:
        _shape_cache: Cache for sorted shapes to avoid redundant sorting.
        total_page: Total number of slides in the presentation.
    """

    def __init__(self):
        super().__init__()
        self._shape_cache = {}

    def __sort_shapes(self, shapes):
        """Sort shapes by vertical position first, then horizontal position.

        Results are cached by the id of the shapes collection for performance.

        Args:
            shapes: A collection of pptx shape objects.

        Returns:
            A list of shapes sorted top-to-bottom, left-to-right.
        """
        cache_key = id(shapes)
        if cache_key not in self._shape_cache:
            self._shape_cache[cache_key] = sorted(
                shapes,
                key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)
            )
        return self._shape_cache[cache_key]

    def __get_bulleted_text(self, paragraph):
        """Extract text from a paragraph, prepending bullet indentation if present.

        Detects bullet markers (character bullets, auto-numbered, and blip bullets)
        and formats the text with appropriate indentation.

        Args:
            paragraph: A pptx paragraph object.

        Returns:
            The paragraph text, optionally prefixed with indentation and bullet marker.
        """
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text

    def __extract(self, shape):
        """Recursively extract text content from a single shape.

        Handles three shape types:
        - Text frames: concatenates paragraph text
        - Tables (shape_type 19): formats as "header: value" pairs per row
        - Group shapes (shape_type 6): recursively extracts from child shapes

        Args:
            shape: A pptx shape object.

        Returns:
            Extracted text string, or empty string if no text found.
        """
        try:
            # First try to get text content from text frame
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)

            # Safely get shape_type
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content directly
                if hasattr(shape, 'text'):
                    return shape.text.strip()
                return ""

            # Handle table shapes (shape_type == 19)
            if shape_type == 19:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append("; ".join([tb.cell(
                        0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
                return "\n".join(rows)

            # Handle group shapes (shape_type == 6) by recursing into children
            if shape_type == 6:
                texts = []
                for p in self.__sort_shapes(shape.shapes):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def __call__(self, fnm, from_page, to_page, callback=None):
        """Parse a PowerPoint file and extract text from each slide.

        Args:
            fnm: File path (string) or binary content (bytes) of the .pptx file.
            from_page: Starting slide index (0-based, inclusive).
            to_page: Ending slide index (0-based, exclusive).
            callback: Optional progress callback function (unused currently).

        Returns:
            A list of strings, one per slide within the page range, containing
            the concatenated text of all shapes on that slide.
        """
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in self.__sort_shapes(slide.shapes):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        return txts
